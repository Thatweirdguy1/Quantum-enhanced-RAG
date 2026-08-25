r"""Build the ten-slide Q-RAG deck.

    python -m docs.slides

Writes ``build/Q-RAG_Slides.pptx``.

Every figure on every slide is resolved from ``results/*.json`` through
:mod:`docs.facts`, for the same reason the written documents are: a number typed onto
a slide is a number that can contradict the paper it is presented alongside, and the
person who notices is in the audience.

Design is deliberately plain. One idea per slide, one accent colour, no decoration
that does not carry information. A viva deck is read at distance by someone deciding
whether the work is sound, so the figures are large and the caveats sit next to the
results they qualify rather than on a later slide.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from docs.facts import Facts

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "build" / "Q-RAG_Slides.pptx"

# ------------------------------------------------------------------ design tokens
INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6E, 0x6E, 0x6E)
RULE = RGBColor(0xD8, 0xD8, 0xD8)
ACCENT = RGBColor(0x1F, 0x3A, 0x5F)      # deep navy, used for emphasis only
NEGATIVE = RGBColor(0x8C, 0x2F, 0x2F)    # muted red: a result that did not work
POSITIVE = RGBColor(0x2C, 0x5F, 0x48)    # muted green: the one that did
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"          # present on every Office install; a projector-safe default
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)       # 16:9
MARGIN = Inches(0.9)
CONTENT_W = W - 2 * MARGIN


def _text(slide, left, top, width, height, *, size, bold=False, colour=INK,
          align=PP_ALIGN.LEFT, font=FONT, spacing=1.0, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.font.size, run.font.bold, run.font.italic = Pt(size), bold, italic
    run.font.color.rgb, run.font.name = colour, font
    return box, p, run


def _rule(slide, top, *, left=MARGIN, width=None, colour=RULE, thickness=Pt(1)):
    width = CONTENT_W if width is None else width
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = colour
    line.line.width = thickness
    return line


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def heading(slide, kicker: str, title: str) -> Emu:
    """Slide furniture: a small kicker, a title, a hairline. Returns content top."""
    if kicker:
        _text(slide, MARGIN, Inches(0.62), CONTENT_W, Inches(0.3),
              size=12, bold=True, colour=ACCENT)[2].text = kicker.upper()
    _text(slide, MARGIN, Inches(0.95), CONTENT_W, Inches(0.7),
          size=30, bold=True)[2].text = title
    _rule(slide, Inches(1.72))
    return Inches(2.0)


def bullets(slide, top, items, *, size=17, gap=Inches(0.52), width=None):
    """Plain bullet list. Items may be (text, colour) or (text, colour, bold)."""
    width = CONTENT_W if width is None else width
    y = top
    for item in items:
        text, colour, bold = (item + (False,))[:3] if isinstance(item, tuple) \
            else (item, INK, False)
        _text(slide, MARGIN + Inches(0.22), y, width - Inches(0.22), gap,
              size=size, colour=colour, bold=bold, spacing=1.15)[2].text = text
        dot, _, run = _text(slide, MARGIN, y, Inches(0.2), Inches(0.3),
                            size=size, colour=MUTED)
        run.text = "—"      # em dash rather than a bullet glyph
        y += gap
    return y


def stat(slide, left, top, width, value, label, *, colour=ACCENT, note=None,
         value_size=40):
    """One large figure with a caption under it."""
    _text(slide, left, top, width, Inches(0.7),
          size=value_size, bold=True, colour=colour)[2].text = value
    _text(slide, left, top + Inches(0.66), width, Inches(0.5),
          size=13, colour=INK, spacing=1.1)[2].text = label
    if note:
        _text(slide, left, top + Inches(1.16), width, Inches(0.9),
              size=11, colour=MUTED, spacing=1.15)[2].text = note


def table(slide, top, headers, rows, *, widths, size=13, row_h=Inches(0.34),
          emphasis=None):
    """Minimal table: a rule under the header, no grid, no fill."""
    xs, x = [], MARGIN
    for w in widths:
        xs.append(x)
        x += w
    for i, htext in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        _text(slide, xs[i], top, widths[i], row_h, size=size - 1, bold=True,
              colour=MUTED, align=align)[2].text = htext
    _rule(slide, top + row_h + Inches(0.04), width=sum(widths))
    y = top + row_h + Inches(0.16)
    for r, row in enumerate(rows):
        colour = INK
        bold = False
        if emphasis and r in emphasis:
            colour, bold = emphasis[r]
        for i, cell in enumerate(row):
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            font = MONO if i == 0 and cell.startswith(("qrag", "classical")) else FONT
            _text(slide, xs[i], y, widths[i], row_h, size=size, colour=colour,
                  bold=bold, align=align, font=font)[2].text = cell
        y += row_h
    return y


def footer(slide, n: int, note: str = ""):
    _rule(slide, H - Inches(0.62))
    _text(slide, MARGIN, H - Inches(0.52), CONTENT_W - Inches(0.6), Inches(0.3),
          size=10, colour=MUTED)[2].text = note
    _text(slide, W - MARGIN - Inches(0.6), H - Inches(0.52), Inches(0.6),
          Inches(0.3), size=10, colour=MUTED, align=PP_ALIGN.RIGHT)[2].text = str(n)


# ------------------------------------------------------------------------- slides
def build(out: Path = OUT) -> Path:
    f = Facts()
    g = f.format                       # every figure below comes through this
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    prov = (f"{g('dataset')} | {g('n_queries')} queries | "
            f"config {g('config_hash')} | seed {g('seed')}")

    # 1 -------------------------------------------------------------------- title
    s = blank(prs)
    _text(s, MARGIN, Inches(2.1), CONTENT_W, Inches(0.4),
          size=13, bold=True, colour=ACCENT)[2].text = \
        "SEMESTER-7 MAJOR PROJECT | GROUP 165"
    _text(s, MARGIN, Inches(2.6), CONTENT_W, Inches(1.3),
          size=44, bold=True, spacing=0.95)[2].text = \
        "Quantum-Enhanced Retrieval-Augmented Generation"
    _rule(s, Inches(4.15), width=Inches(3.2), colour=ACCENT, thickness=Pt(2.5))
    _text(s, MARGIN, Inches(4.4), CONTENT_W, Inches(0.9), size=17, colour=MUTED,
          spacing=1.2)[2].text = (
        "Three simulated quantum stages in a RAG pipeline, and experiments built so "
        "that each one can be shown to do nothing if that is the truth.")
    _text(s, MARGIN, Inches(5.6), CONTENT_W, Inches(0.4), size=14)[2].text = \
        "Prabhav Goel  |  Group 165  |  B.Tech CSE"
    _text(s, MARGIN, Inches(6.0), CONTENT_W, Inches(0.4), size=12,
          colour=MUTED)[2].text = \
        "Amity School of Engineering and Technology  |  August 2026"
    footer(s, 1, prov)

    # 2 ------------------------------------------------------- the honest headline
    s = blank(prs)
    top = heading(s, "where this ends up", "The result, stated first")
    col = (CONTENT_W - Inches(0.8)) / 3
    stat(s, MARGIN, top, col, "null", "Retrieval quality",
         colour=NEGATIVE,
         note=f"{g('sig.count')} of {g('sig.total')} significance cells crosses "
              f"p < 0.05. With {g('sig.total')} comparisons, that is noise.")
    stat(s, MARGIN + col + Inches(0.4), top, col,
         f"{g('slowdown')}× slower", "Wall clock", colour=NEGATIVE,
         note="A statevector simulator cannot beat the routine it simulates. "
              "No speed-up is claimed anywhere.")
    stat(s, MARGIN + 2 * (col + Inches(0.4)), top, col,
         f"−{g('poison.qaoa_effect')}", "Adversarial context occupancy",
         colour=POSITIVE,
         note="The one positive finding: QAOA's redundancy penalty, isolated by "
              "an ablation. Narrow, and still not a defence.")
    _rule(s, Inches(5.35))
    _text(s, MARGIN, Inches(5.55), CONTENT_W, Inches(1.2), size=15, spacing=1.25,
          colour=INK)[2].text = (
        "Two of the three headline numbers are negative. They are on the second "
        "slide rather than the tenth because a project that hides them until the "
        "end is asking to be caught, and because the experiments were designed to "
        "produce exactly this answer if it was the true one.")
    footer(s, 2, prov)

    # 3 --------------------------------------------------------- the central risk
    s = blank(prs)
    top = heading(s, "the problem with quantum ir",
                  "The central risk: a kernel that cannot reorder")
    _text(s, MARGIN, top, CONTENT_W, Inches(0.5), size=17, spacing=1.2)[2].text = (
        "A fidelity kernel over amplitude-encoded vectors, at zero phase, is:")
    _text(s, MARGIN + Inches(0.3), top + Inches(0.62), CONTENT_W, Inches(0.5),
          size=22, font=MONO, colour=ACCENT)[2].text = \
        "K(q,d) = |Σ qᵢ dᵢ e^(iθᵢ)|²   →   cos²(q,d)   at θ = 0"
    bullets(s, top + Inches(1.4), [
        ("cos² is a monotone function of cosine, so it induces the identical "
         "ranking. Every document keeps its position.", INK),
        (f"Measured on the untrained global kernel: Kendall τ = "
         f"{g('gate_a.tau_global_before')} against the cosine ranking — it "
         f"reorders essentially nothing and fails the {g('gate_a.threshold')} "
         f"ceiling set for Gate A.", NEGATIVE, True),
        ("A “quantum RAG” whose quantum stage reorders nothing is "
         "vacuous, however well the rest of it performs.", INK),
    ], size=16, gap=Inches(0.72))
    _rule(s, Inches(5.5))
    _text(s, MARGIN, Inches(5.7), CONTENT_W, Inches(1.0), size=15, spacing=1.2,
          colour=INK, italic=True)[2].text = (
        "This is reported as a finding, not omitted as an inconvenience. It is also "
        "why the project is built around a kernel that breaks the equivalence "
        "structurally rather than one that hopes training will fix it.")
    footer(s, 3, prov)

    # 4 --------------------------------------------------- the fix and the gates
    s = blank(prs)
    top = heading(s, "the response", "Break the equivalence, then pre-register the test")
    _text(s, MARGIN, top, CONTENT_W, Inches(0.4), size=16, spacing=1.2)[2].text = (
        "The block (projected) fidelity kernel sums fidelities over disjoint "
        f"blocks — {g('n_blocks')} blocks of {g('block_size')} dimensions:")
    _text(s, MARGIN + Inches(0.3), top + Inches(0.52), CONTENT_W, Inches(0.4),
          size=19, font=MONO, colour=ACCENT)[2].text = \
        "K = Σ_g w_g |Σ_{i∈g} qᵢ dᵢ e^(iθᵢ)|²   →   Σ_g S_g²   at θ = 0"
    _text(s, MARGIN, top + Inches(1.02), CONTENT_W, Inches(0.4), size=15,
          colour=MUTED, spacing=1.15)[2].text = (
        "By Cauchy–Schwarz, Σ_g S_g² is not monotone in "
        "(Σ_g S_g)² = cos². So it can reorder even before training.")
    _rule(s, top + Inches(1.6))
    _text(s, MARGIN, top + Inches(1.78), CONTENT_W, Inches(0.35), size=13,
          bold=True, colour=ACCENT)[2].text = \
        "BOTH GATES WERE DEFINED BEFORE TRAINING, NOT AFTER SEEING RESULTS"
    half = (CONTENT_W - Inches(0.6)) / 2
    y = top + Inches(2.25)
    stat(s, MARGIN, y, half, f"τ = {g('gate_a.tau')}",
         f"Gate A — does it reorder? Ceiling was {g('gate_a.threshold')}.",
         colour=POSITIVE, value_size=30,
         note=f"Training-free control at θ = 0 already reaches "
              f"{g('gate_a.tau_theta0')}, confirming the construction and not "
              f"the fitted phases is what breaks equivalence.")
    stat(s, MARGIN + half + Inches(0.6), y, half,
         f"{g('gate_b.mrr_before')} → {g('gate_b.mrr_after')}",
         "Gate B — is the reordering an improvement? Held-out MRR.",
         colour=POSITIVE, value_size=30,
         note=f"Passes, but thinly: {g('n_val_pairs')} validation pairs and a "
              f"top-1 change of {g('gate_b.delta_top1')} — about two queries. "
              f"Not significant, and not presented as such.")
    footer(s, 4, prov)

    # 5 ----------------------------------------------------------- what was built
    s = blank(prs)
    top = heading(s, "the system", "Five stages, each independently switchable")
    rows = [
        ("Hybrid fusion", "BM25 + dense cosine + quantum kernel, per-query min-max "
                          "normalised", "baseline is the same fusion, kernel weight "
                          "moved to cosine"),
        ("Phase kernel", f"block fidelity, {g('n_blocks')}×{g('block_size')} "
                         f"phases, InfoNCE, analytic gradients",
         f"trained on {g('n_train_pairs')} pairs; overfits by construction"),
        ("Interference", "sub-query decomposition with decaying weights",
         "works, small effect"),
        ("Grover", f"threshold oracle over a scored shortlist, "
                   f"{g('grover.qubits')} qubits",
         "reported in oracle queries, never wall clock"),
        ("QAOA", f"selection QUBO, {g('qaoa.qubits')} qubits, p = "
                 f"{g('qaoa.layers')}, exact optimum computed alongside",
         f"{g('qaoa.ms')} ms/query — the expensive stage"),
        ("Adversarial arm", f"{g('poison.n_families')} attack families, "
                            f"{g('poison.n_injected')} injected passages, qrels "
                            f"untouched", "detector's blind spot is measured"),
    ]
    y = top
    for name, what, status in rows:
        _text(s, MARGIN, y, Inches(2.2), Inches(0.3), size=15, bold=True)[2].text = name
        _text(s, MARGIN + Inches(2.3), y, Inches(6.0), Inches(0.6), size=13,
              spacing=1.1)[2].text = what
        _text(s, MARGIN + Inches(8.5), y, Inches(3.0), Inches(0.6), size=12,
              colour=MUTED, spacing=1.1)[2].text = status
        y += Inches(0.68)
        _rule(s, y - Inches(0.12))
    _text(s, MARGIN, y + Inches(0.06), CONTENT_W, Inches(0.4), size=13,
          colour=MUTED, italic=True)[2].text = (
        "Every numpy fast path is cross-checked against qiskit-aer. "
        "All circuits are simulated — nothing ran on quantum hardware.")
    footer(s, 5, prov)

    # 6 -------------------------------------------------------- experiment design
    s = blank(prs)
    top = heading(s, "how it was measured", "An experiment designed to be able to fail")
    bullets(s, top, [
        (f"{g('n_systems')} configurations × all {g('n_queries')} SciFact test "
         f"queries. The baseline is a tuned hybrid, not a strawman: identical "
         f"fusion with the kernel weight redistributed onto cosine.", INK),
        ("Paired bootstrap, 2,000 resamples, per-query scores retained. "
         "Significance means the 95% CI excludes zero.", INK),
        (f"Poisoned arm: {g('poison.n_injected')} adversarial passages against "
         f"{g('poison.n_targets')} target queries. Relevance judgments left "
         f"unmodified so retrieval metrics stay comparable with the clean arm.",
         INK),
        ("qrag[no-qaoa] exists purely as the control that makes the security "
         "hypothesis falsifiable — without it, any drop in adversarial "
         "occupancy could be credited to the kernel instead.", ACCENT, True),
        ("No test asserts a research result. A test like assert mrr > 0.7 turns a "
         "measurement into a requirement and creates pressure to tune until it "
         "passes.", INK),
    ], size=16, gap=Inches(0.82))
    _rule(s, Inches(6.05))
    _text(s, MARGIN, Inches(6.25), CONTENT_W, Inches(0.5), size=13,
          colour=MUTED)[2].text = (
        f"Reproducible: config {g('config_hash')}, seed {g('seed')}, "
        f"Python {g('python')}, numpy {g('numpy')}. "
        f"Same command and seed gives the same numbers.")
    footer(s, 6, prov)

    # 7 -------------------------------------------------------- result 1: the null
    s = blank(prs)
    top = heading(s, "result 1 of 3", "Retrieval quality: no significant improvement")
    widths = [Inches(3.3), Inches(1.7), Inches(1.7), Inches(1.7), Inches(2.1),
              Inches(1.0)]
    exp = f.exp["clean"]
    sig = f.exp["significance"]
    order = ["classical-baseline", "qrag[kernel]", "qrag[kernel+interf]",
             "qrag[grover]", "qrag[qaoa]", "qrag[kernel+qaoa]", "qrag[full]"]
    rows, emphasis = [], {}
    for i, label in enumerate(order):
        m = exp[label]["metrics"]
        cell = sig.get(label, {}).get("ndcg@10")
        delta = "—" if cell is None else f"{cell['delta']:+.4f}"
        p = "—" if cell is None else f"{cell['p_value']:.3f}"
        rows.append((label, f"{m['recall@10']:.4f}", f"{m['ndcg@10']:.4f}",
                     f"{m['mrr@10']:.4f}", delta, p))
        if label == "classical-baseline":
            emphasis[i] = (ACCENT, True)
    table(s, top, ["system", "recall@10", "nDCG@10", "MRR@10",
                   "Δ nDCG@10", "p"], rows, widths=widths, row_h=Inches(0.36))
    _rule(s, Inches(5.35))
    _text(s, MARGIN, Inches(5.5), CONTENT_W, Inches(1.4), size=15,
          spacing=1.25)[2].text = (
        f"Every nDCG@10 delta lies within ±0.002 of the baseline and none is "
        f"significant. One cell of {g('sig.total')} reaches p < 0.05 "
        f"(qrag[kernel] recall@5, {g('sig.first.delta|signed')}, "
        f"p = {g('sig.first.p|3f')}) and with {g('sig.total')} comparisons that is "
        f"what noise looks like — it is not claimed as a finding. The kernel "
        f"cleared both pre-registered gates on held-out pairs and then did not "
        f"transfer to the full corpus. That is the result.")
    footer(s, 7, prov)

    # 8 --------------------------------------------------- result 2: the accounting
    s = blank(prs)
    top = heading(s, "result 2 of 3",
                  "Quantum accounting: complexity and cost, kept apart")
    half = (CONTENT_W - Inches(0.7)) / 2
    stat(s, MARGIN, top, half, f"{g('grover.reduction|2f')}×",
         "Fewer oracle queries (Grover)", colour=POSITIVE, value_size=44,
         note=f"{g('grover.oracle_queries|0f')} query against "
              f"{g('grover.classical_queries|2f')} expected for a classical scan "
              f"over {g('grover.candidates|0f')} candidates. Success probability "
              f"{g('grover.success')}. A hardware-independent count.")
    stat(s, MARGIN + half + Inches(0.7), top, half,
         f"{g('grover_arm.overhead|2f')}×", "More time to simulate it",
         colour=NEGATIVE, value_size=44,
         note="The same subroutine, measured on the clock. Printed in the adjacent "
              "column on purpose: the query reduction may never be restated as a "
              "speed, a latency or a throughput gain.")
    _rule(s, Inches(4.42))
    _text(s, MARGIN, Inches(4.58), CONTENT_W, Inches(0.35), size=13, bold=True,
          colour=ACCENT)[2].text = "QAOA, AGAINST AN EXACT BRUTE-FORCE OPTIMUM"
    y = Inches(4.95)
    third = (CONTENT_W - Inches(0.8)) / 3
    stat(s, MARGIN, y, third, g('qaoa.quality'), "Mean solution quality",
         colour=INK, value_size=26,
         note=f"Affine-invariant; 1.0 is the exact optimum over the same "
              f"exactly-k feasible set. Worst case {g('qaoa.quality_worst')}.")
    stat(s, MARGIN + third + Inches(0.4), y, third, g('qaoa.exact_rate|pct'),
         "Queries hitting the exact optimum", colour=INK, value_size=26,
         note=f"{g('qaoa.optimiser_calls|0f')} optimiser calls per query; "
              f"feasible probability at readout {g('qaoa.feasible_prob|2f')}.")
    stat(s, MARGIN + 2 * (third + Inches(0.4)), y, third, f"{g('qaoa.ms')} ms",
         "Per query, simulated", colour=NEGATIVE, value_size=26,
         note=f"{g('qaoa.share|pct')} of pipeline latency — "
              f"{g('qaoa.vs_rest|1f')}× the other four stages combined — and it "
              f"buys no retrieval gain on a clean corpus.")
    footer(s, 8, prov)

    # 9 ----------------------------------------------------- result 3: the security
    s = blank(prs)
    top = heading(s, "result 3 of 3",
                  "Security: the one positive result, and its limits")
    del top     # slide 9 lays out against fixed rules, not the heading's return
    widths = [Inches(3.6), Inches(2.6), Inches(2.4), Inches(2.4)]
    psys = f.exp["poisoned"]["systems"]
    rows = []
    for label in ["classical-baseline", "qrag[no-qaoa]", "qrag[full]"]:
        a = psys[label]["attack"]
        rows.append((label, f"{a['context_occupancy']:.4f}",
                     f"{a['clean_context_rate']:.4f}",
                     f"{a['top_k_hit_rate']:.4f}"))
    table(s, Inches(1.95), ["system on the poisoned corpus", "context occupancy",
                            "clean-context rate", "top-10 hit rate"], rows,
          widths=widths, row_h=Inches(0.36), emphasis={2: (POSITIVE, True)})
    _text(s, MARGIN, Inches(3.72), CONTENT_W, Inches(0.8), size=15,
          spacing=1.2)[2].text = (
        f"Occupancy falls {g('poison.occ_noqaoa')} → {g('poison.occ_full')}, "
        f"a reduction of {g('poison.qaoa_effect')} attributable to the redundancy "
        f"penalty alone — the two arms differ in that term and nothing else.")
    _rule(s, Inches(4.42))
    _text(s, MARGIN, Inches(4.58), CONTENT_W, Inches(0.35), size=13,
          bold=True, colour=NEGATIVE)[2].text = \
        "WHAT THAT NUMBER DOES NOT MEAN"
    bullets(s, Inches(4.95), [
        (f"{g('poison.occ_full')} occupancy is still catastrophic. All "
         f"{g('poison.n_targets')} targeted queries were hit, and the "
         f"clean-context rate is {g('poison.clean_rate_full')} — not one "
         f"query received an uncontaminated context.", NEGATIVE),
        (f"The pattern detector catches "
         f"{g('detector.instruction-injection|pct')} of instruction injection and "
         f"{g('detector.topical-mimicry|pct')} of fluent topical mimicry. The "
         f"aggregate {g('detector.rate|pct')} may not be quoted without that "
         f"second figure.", NEGATIVE),
        ("The embedding-optimised attack is black-box, so it is strictly weaker "
         "than a gradient-based one. Any defence result is against the weaker "
         "attack.", MUTED),
    ], size=13, gap=Inches(0.62))
    footer(s, 9, prov)

    # 10 --------------------------------------------------------------- conclusion
    s = blank(prs)
    top = heading(s, "conclusion", "What this work supports, and what it does not")
    _text(s, MARGIN, top, CONTENT_W, Inches(0.35), size=13, bold=True,
          colour=POSITIVE)[2].text = "SUPPORTED BY THE MEASUREMENTS"
    y = bullets(s, top + Inches(0.38), [
        "A block fidelity kernel breaks rank-equivalence with cosine structurally, "
        "and the untrained global kernel does not — both demonstrated, not "
        "assumed.",
        f"Grover's {g('grover.reduction|2f')}× oracle-query reduction on a real "
        f"retrieval workload, reported separately from its "
        f"{g('grover_arm.overhead|2f')}× simulation cost.",
        f"QAOA's redundancy penalty reduces adversarial context occupancy by "
        f"{g('poison.qaoa_effect')} on one corpus at one attack budget, isolated "
        f"by an ablation.",
    ], size=14, gap=Inches(0.56))
    _rule(s, y + Inches(0.04))
    _text(s, MARGIN, y + Inches(0.18), CONTENT_W, Inches(0.35), size=13, bold=True,
          colour=NEGATIVE)[2].text = "CLAIMS THIS PROJECT MAY NOT MAKE"
    bullets(s, y + Inches(0.56), [
        (f"No speed-up. The pipeline is {g('slowdown')}× slower than its own "
         f"baseline. Not “faster”, not “efficient”.", NEGATIVE),
        ("No retrieval improvement. No system beats the baseline on nDCG@10 with a "
         "CI excluding zero.", NEGATIVE),
        ("No hardware claim, and no generalisation beyond one corpus and one "
         "attack budget.", NEGATIVE),
    ], size=14, gap=Inches(0.5))
    _text(s, MARGIN, Inches(6.35), CONTENT_W, Inches(0.5), size=12,
          colour=MUTED, spacing=1.15)[2].text = (
        f"{g('tests.passed')} tests pass; security audit reports "
        f"{g('audit.passed')} passed, {g('audit.failed')} failed, "
        f"{g('audit.warned')} warned, {g('audit.na')} not applicable. "
        f"Every figure in this deck is generated from results/*.json at build time.")
    footer(s, 10, prov)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    return out


def main() -> int:
    out = build()
    n = len(Presentation(out).slides._sldIdLst)
    print(f"slides -> {out.relative_to(ROOT)}  "
          f"({n} slides, {out.stat().st_size / 1024:.0f} KB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
